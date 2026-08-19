#!/usr/bin/env python3
"""Generate the presentation-ready PSTMO slide deck from audited repository evidence."""

from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDE_DIR = ROOT / "slide"
ASSETS = ROOT / "docs" / "pstmo_bao_cao_toan_dien_assets"
VISUALS_3D = ASSETS / "visuals_3d"
OUTPUT = SLIDE_DIR / "slide.pptx"

FONT = "Noto Sans"
WIDE = 13.333
HIGH = 7.5

NAVY = "102A43"
BLUE = "1D6FB8"
TEAL = "0F766E"
GREEN = "169C52"
ORANGE = "E8871E"
RED = "C53B3C"
PURPLE = "7257B5"
INK = "17212B"
MUTED = "5B6876"
LIGHT = "F4F7FA"
PANEL = "EAF1F7"
LINE = "D6DEE7"
WHITE = "FFFFFF"
PALE_GREEN = "E7F6ED"
PALE_ORANGE = "FFF3E5"
PALE_RED = "FBEAEC"
PALE_BLUE = "E7F1FA"
PALE_PURPLE = "F0ECFA"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_cell_border(cell, color: str = LINE, width: float = 0.8) -> None:
    """Apply borders through the underlying DrawingML table cell."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls

    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tag = edge.split(":")[1]
        old = tc_pr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if old is not None:
            tc_pr.remove(old)
        tc_pr.append(
            parse_xml(
                f'<{edge} {nsdecls("a")} w="{int(width * 12700)}">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                '<a:prstDash val="solid"/>'
                '</' + edge + '>'
            )
        )


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font: str = FONT,
    italic: bool = False,
    line_spacing: float | None = 1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    runs: Sequence[dict],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for item in runs:
        run = p.add_run()
        run.text = item["text"]
        run.font.name = item.get("font", FONT)
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.color.rgb = rgb(item.get("color", color))
    return box


def add_bullets(
    slide,
    items: Iterable[str | tuple[str, int]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 17,
    color: str = INK,
    accent: str = BLUE,
    spacing: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for idx, raw in enumerate(items):
        text, level = raw if isinstance(raw, tuple) else (raw, 0)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ("• " if level == 0 else "   – ") + text
        p.level = level
        p.font.name = FONT
        p.font.size = Pt(size - level * 1.5)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        if p.runs:
            p.runs[0].font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str | None = LINE,
    radius: bool = True,
    line_width: float = 1.0,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_width)
    return shp


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = LINE, width: float = 1.5):
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(width)
    return ln


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = BLUE, width: float = 2.0):
    ln = add_line(slide, x1, y1, x2, y2, color=color, width=width)
    return ln


def add_pill(slide, text: str, x: float, y: float, w: float, *, fill: str, color: str = WHITE, size: float = 11):
    add_rect(slide, x, y, w, 0.32, fill=fill, line=None)
    add_text(slide, text, x + 0.04, y + 0.01, w - 0.08, 0.26, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_title(slide, title: str, section: str, page: int, *, subtitle: str | None = None):
    add_text(slide, title, 0.62, 0.30, 10.9, 0.55, size=25, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, 0.65, 0.86, 11.3, 0.30, size=11.5, color=MUTED)
    add_pill(slide, section.upper(), 11.35, 0.34, 1.35, fill=PALE_BLUE, color=BLUE, size=9.2)
    add_line(slide, 0.62, 1.18, 12.72, 1.18, color=LINE, width=1.0)
    add_text(slide, f"{page:02d}", 12.20, 7.05, 0.50, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_caption(slide, text: str, x: float, y: float, w: float):
    add_text(slide, text, x, y, w, 0.30, size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


def add_card(slide, x: float, y: float, w: float, h: float, *, title: str, body: str, accent: str = BLUE, fill: str = WHITE, title_size: float = 15, body_size: float = 12.5):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE)
    add_rect(slide, x, y, 0.08, h, fill=accent, line=None, radius=False)
    add_text(slide, title, x + 0.24, y + 0.15, w - 0.40, 0.35, size=title_size, color=NAVY, bold=True)
    add_text(slide, body, x + 0.24, y + 0.58, w - 0.40, h - 0.70, size=body_size, color=MUTED)


def add_metric(slide, x: float, y: float, w: float, h: float, *, value: str, label: str, accent: str = GREEN, note: str | None = None):
    add_rect(slide, x, y, w, h, fill=WHITE, line=LINE)
    add_text(slide, value, x + 0.14, y + 0.12, w - 0.28, 0.50, size=23, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.14, y + 0.66, w - 0.28, 0.38, size=11.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + 0.14, y + 1.05, w - 0.28, h - 1.13, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_equation(slide, text: str, x: float, y: float, w: float, h: float, *, color: str = NAVY, fill: str = LIGHT, size: float = 20):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE)
    add_text(slide, text, x + 0.12, y + 0.05, w - 0.24, h - 0.10, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, *, background: str = WHITE):
    add_rect(slide, x, y, w, h, fill=background, line=LINE, radius=False)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    pic = slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))
    return pic


def add_picture_cover(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as im:
        iw, ih = im.size
    frame_ratio = w / h
    image_ratio = iw / ih
    pic = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if image_ratio > frame_ratio:
        crop = (1.0 - frame_ratio / image_ratio) / 2.0
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1.0 - image_ratio / frame_ratio) / 2.0
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def add_section_slide(prs, number: str, title: str, subtitle: str, accent: str, image_path: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_text(slide, number, 0.72, 0.76, 1.40, 1.10, size=54, color=accent, bold=True)
    add_text(slide, title, 0.76, 2.02, 6.25, 1.38, size=31, color=WHITE, bold=True)
    add_text(slide, subtitle, 0.80, 3.58, 5.95, 1.20, size=16, color="C8D5E3")
    add_rect(slide, 0.80, 5.30, 1.05, 0.08, fill=accent, line=None, radius=False)
    if image_path:
        add_picture_cover(slide, image_path, 7.36, 0.0, 5.97, 7.50)
        overlay = add_rect(slide, 7.05, 0.0, 0.95, 7.5, fill=NAVY, line=None, radius=False)
        overlay.fill.transparency = 30
    return slide


def add_node(slide, text: str, x: float, y: float, w: float, h: float, *, fill: str, color: str = WHITE, size: float = 12):
    add_rect(slide, x, y, w, h, fill=fill, line=None)
    add_text(slide, text, x + 0.06, y + 0.03, w - 0.12, h - 0.06, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_circle_label(slide, text: str, x: float, y: float, d: float, *, fill: str = BLUE, color: str = WHITE, size: float = 13):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    add_text(slide, text, x, y, d, d, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_table(slide, data: Sequence[Sequence[str]], x: float, y: float, w: float, h: float, *, col_widths: Sequence[float] | None = None, header_fill: str = NAVY, font_size: float = 11.5):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(header_fill if r == 0 else (WHITE if r % 2 else LIGHT))
            set_cell_border(cell)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c != 0 or r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size if r else font_size - 0.2)
                    run.font.bold = r == 0
                    run.font.color.rgb = rgb(WHITE if r == 0 else INK)
    return table


def build() -> None:
    geo = json.loads((ASSETS / "benchmark_hinh_hoc_tong_hop.json").read_text())
    execution = json.loads((ASSETS / "execution_aggregate_5planners_7env.json").read_text())

    prs = Presentation()
    prs.slide_width = Inches(WIDE)
    prs.slide_height = Inches(HIGH)
    props = prs.core_properties
    props.title = "PSTMO – Làm mượt đường đi và tối ưu hóa thao tác chuyển hướng"
    props.subject = "ROS 2 Navigation2, robot vi sai, Bézier bậc năm G², RViz2, Gazebo"
    props.author = "NGUYỄN TIẾN CƯƠNG"
    props.keywords = "PSTMO, ROS 2, Nav2, path smoothing, differential drive, RViz2, Gazebo"
    props.comments = "Generated from audited repository figures and benchmark JSON."

    # 01 — Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_picture_contain(slide, VISUALS_3D / "robot_isometric_clean.png", 7.62, 0.0, 5.71, 7.50, background=WHITE)
    add_rect(slide, 7.15, 0.0, 1.00, 7.5, fill=WHITE, line=None, radius=False)
    add_pill(slide, "ROS 2 NAVIGATION2 · ROBOT VI SAI", 0.72, 0.62, 3.35, fill=PALE_BLUE, color=BLUE, size=10.5)
    add_text(slide, "PSTMO", 0.74, 1.38, 5.95, 0.72, size=34, color=NAVY, bold=True)
    add_text(slide, "Làm mượt đường đi và tối ưu hóa\nthao tác chuyển hướng", 0.74, 2.18, 6.25, 1.40, size=27, color=INK, bold=True)
    add_text(slide, "Path smoothing and turning-maneuver optimization", 0.78, 3.78, 5.95, 0.48, size=15.5, color=TEAL, italic=True)
    add_line(slide, 0.78, 4.54, 5.95, 4.54, color=LINE, width=1.2)
    add_text(slide, "NGUYỄN TIẾN CƯƠNG", 0.78, 4.82, 4.80, 0.38, size=15, color=NAVY, bold=True)
    add_text(slide, "Bằng chứng: mô phỏng Gazebo Harmonic · trực quan RViz2 · 175 lượt thực thi", 0.78, 5.31, 5.95, 0.70, size=12.5, color=MUTED)
    add_text(slide, "2026", 0.78, 6.70, 1.00, 0.30, size=11, color=MUTED)
    add_text(slide, "Ảnh bìa: dựng trực tiếp từ STL/SDF", 8.28, 7.03, 4.62, 0.22, size=8.5, color=MUTED, italic=True, align=PP_ALIGN.RIGHT)

    # 02 — Research question
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Câu hỏi nghiên cứu", "Mở đầu", 2, subtitle="Từ đường đi tồn tại trên bản đồ tới chuyển động có thể thực thi")
    add_text(slide, "Làm thế nào thay các góc gãy của đường do bộ lập kế hoạch tạo bằng thao tác chuyển hướng…", 0.78, 1.55, 11.75, 0.58, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    cards = [
        ("MƯỢT HÌNH HỌC", "Tiếp tuyến và độ cong không nhảy tại điểm nối.", BLUE),
        ("KHẢ THI ĐỘNG HỌC", "Không vượt giới hạn thân xe, bánh xe và gia tốc.", PURPLE),
        ("AN TOÀN", "Toàn bộ hình bao robot không quét vào vật cản.", GREEN),
        ("CÓ LỢI VỀ THỜI GIAN", "Chỉ tạo đoạn chuyển tiếp khi nhanh hơn quay tại chỗ.", ORANGE),
    ]
    for i, (title, body, accent) in enumerate(cards):
        add_card(slide, 0.78 + i * 3.05, 2.62, 2.78, 2.25, title=title, body=body, accent=accent, title_size=13.2, body_size=12.2)
    add_equation(slide, "Đầu vào: nav_msgs/Path thô  →  Đầu ra: nav_msgs/Path đã xử lý", 2.02, 5.42, 9.30, 0.72, color=TEAL, fill=PALE_GREEN, size=18)
    add_text(slide, "PSTMO tối ưu đường hình học; bộ điều khiển bám đường vẫn là khối phát lệnh vận tốc v, ω.", 1.38, 6.39, 10.50, 0.36, size=13.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 03 — Contributions
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Đóng góp và phạm vi khẳng định", "Mở đầu", 3)
    items = [
        ("1", "Đoạn chuyển tiếp Bézier bậc năm G²", "Nối đoạn thẳng–đường cong với κ=0 tại hai đầu."),
        ("2", "Chọn trạng thái theo nhiều điều kiện", "Hình học → động học → vùng quét hình bao → thời gian."),
        ("3", "Tối ưu toàn đường bằng quy hoạch động", "Ngăn vùng cắt của hai góc kề nhau chồng lấn."),
        ("4", "Kiểm chứng trên đúng hệ thống mô phỏng", "7 môi trường × 5 bộ lập kế hoạch × 5 phương án."),
    ]
    for i, (n, title, body) in enumerate(items):
        y = 1.48 + i * 1.22
        add_circle_label(slide, n, 0.90, y + 0.05, 0.52, fill=[BLUE, TEAL, PURPLE, GREEN][i])
        add_text(slide, title, 1.66, y, 4.65, 0.34, size=16, color=NAVY, bold=True)
        add_text(slide, body, 1.66, y + 0.42, 4.95, 0.48, size=12.5, color=MUTED)
    add_rect(slide, 7.05, 1.48, 5.15, 4.95, fill=LIGHT, line=LINE)
    add_text(slide, "PSTMO chứng minh điều gì?", 7.39, 1.76, 4.47, 0.38, size=18, color=NAVY, bold=True)
    add_bullets(slide, [
        "Đường đầu ra đạt các điều kiện hình học, động học và hình bao robot trên bản đồ chi phí dùng khi làm mượt.",
        "Chỉ số hình học và thời gian mô phỏng được đo trên dữ liệu ghép cặp.",
        "Không đồng nghĩa với bảo đảm tuyệt đối trên robot thật hoặc trước vật cản động.",
    ], 7.42, 2.35, 4.34, 2.45, size=14, spacing=12)
    add_pill(slide, "KẾT LUẬN CÓ ĐIỀU KIỆN", 8.00, 5.40, 3.24, fill=PALE_ORANGE, color=ORANGE, size=10)

    # 04 — Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Mạch trình bày", "Mở đầu", 4, subtitle="Mỗi phần trả lời một câu hỏi cụ thể")
    agenda = [
        ("I", "HỆ THỐNG & BÀI TOÁN", "PSTMO nằm ở đâu trong Nav2? Vì sao góc gãy khó chạy?", BLUE),
        ("II", "PHƯƠNG PHÁP PSTMO", "Đường cong được sinh, kiểm tra và chọn như thế nào?", TEAL),
        ("III", "THỰC NGHIỆM & KẾT QUẢ", "Mô hình gồm gì, thử ra sao và kết quả nói lên điều gì?", GREEN),
    ]
    for i, (n, title, body, accent) in enumerate(agenda):
        y = 1.55 + i * 1.62
        add_rect(slide, 0.90, y, 11.55, 1.20, fill=WHITE, line=LINE)
        add_rect(slide, 0.90, y, 1.12, 1.20, fill=accent, line=None, radius=False)
        add_text(slide, n, 0.90, y, 1.12, 1.20, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, 2.30, y + 0.17, 3.55, 0.32, size=16.5, color=NAVY, bold=True)
        add_text(slide, body, 5.82, y + 0.17, 6.10, 0.65, size=13.2, color=MUTED)

    # 05 — Nav2 architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Nav2 biến bản đồ thành chuyển động như thế nào?", "Hệ thống", 5)
    nodes = [
        ("Bản đồ + cảm biến\nMap + sensors", BLUE),
        ("Định vị\nLocalization", PURPLE),
        ("Lập kế hoạch toàn cục\nGlobal planner", BLUE),
        ("Làm mượt đường\nPath smoother", TEAL),
        ("Điều khiển bám đường\nPath-following controller", ORANGE),
        ("Robot", GREEN),
    ]
    xs = [0.55, 2.65, 4.64, 6.73, 8.82, 11.20]
    widths = [1.62, 1.50, 1.62, 1.62, 1.88, 1.30]
    for i, ((label, accent), x, w) in enumerate(zip(nodes, xs, widths)):
        add_node(slide, label, x, 2.18, w, 0.86, fill=accent, size=10.5)
        if i < len(nodes) - 1:
            add_arrow(slide, x + w + 0.06, 2.61, xs[i + 1] - 0.06, 2.61, color=LINE, width=1.6)
    add_card(slide, 0.72, 3.65, 3.67, 1.60, title="Bộ lập kế hoạch (planner)", body="Trả lời: đi qua đâu?\nĐầu ra: chuỗi tư thế nav_msgs/Path.", accent=BLUE, title_size=14)
    add_card(slide, 4.82, 3.65, 3.67, 1.60, title="Bộ làm mượt (smoother)", body="Trả lời: chỉnh hình học đường ra sao?\nPSTMO nằm ở đây.", accent=TEAL, title_size=14)
    add_card(slide, 8.94, 3.65, 3.67, 1.60, title="Bộ điều khiển (controller)", body="Trả lời: ngay lúc này phát v, ω bao nhiêu?\nĐầu ra: cmd_vel.", accent=ORANGE, title_size=14)
    add_equation(slide, "Trong thí nghiệm: FollowPath nhận đường sau làm mượt; đường gốc (Raw) là nhánh đối chứng.", 1.35, 5.84, 10.65, 0.72, color=NAVY, fill=LIGHT, size=15.2)

    # 06 — Path vs trajectory
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Phân biệt “đường đi” và “quỹ đạo”", "Hệ thống", 6, subtitle="Dùng đúng thuật ngữ để tránh nhầm trách nhiệm của thuật toán")
    add_rect(slide, 0.78, 1.50, 5.78, 4.75, fill=PALE_BLUE, line=LINE)
    add_text(slide, "PATH · ĐƯỜNG HÌNH HỌC", 1.08, 1.80, 5.18, 0.42, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "P = {p₀, p₁, …, pₙ}\npᵢ = (xᵢ, yᵢ, ψᵢ)", 1.42, 2.48, 4.48, 1.12, color=NAVY, fill=WHITE, size=20)
    add_bullets(slide, ["Mô tả vị trí và hướng đi qua.", "Không chứa thời điểm, vận tốc, gia tốc."], 1.28, 4.05, 4.85, 1.25, size=14)
    add_rect(slide, 6.78, 1.50, 5.78, 4.75, fill=PALE_GREEN, line=LINE)
    add_text(slide, "TRAJECTORY · QUỸ ĐẠO THEO THỜI GIAN", 7.08, 1.80, 5.18, 0.42, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "q(t) = [x(t), y(t), ψ(t)]\n+ v(t), ω(t), a(t)", 7.42, 2.48, 4.48, 1.12, color=NAVY, fill=WHITE, size=20)
    add_bullets(slide, ["Gắn trạng thái với thời gian thực thi.", "Cho phép kiểm tra vận tốc và gia tốc."], 7.28, 4.05, 4.85, 1.25, size=14)
    add_text(slide, "PSTMO trả về đường hình học, nhưng dùng mô hình thời gian để loại đường cong không khả thi.", 1.50, 6.48, 10.40, 0.36, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 07 — Corner problem
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Vì sao góc gãy khó thực thi?", "Hệ thống", 7)
    add_rect(slide, 0.72, 1.55, 5.35, 4.95, fill=LIGHT, line=LINE)
    # Geometry sketch
    add_line(slide, 1.25, 5.35, 3.45, 5.35, color=BLUE, width=5)
    add_line(slide, 3.45, 5.35, 3.45, 2.15, color=BLUE, width=5)
    add_circle_label(slide, "V", 3.23, 5.13, 0.44, fill=RED, size=11)
    add_text(slide, "hướng vào", 1.47, 5.62, 1.25, 0.28, size=11, color=MUTED)
    add_text(slide, "hướng ra", 3.68, 2.20, 1.18, 0.28, size=11, color=MUTED)
    add_text(slide, "θ", 3.00, 4.55, 0.38, 0.38, size=22, color=ORANGE, bold=True)
    add_arrow(slide, 2.22, 4.56, 3.13, 4.56, color=ORANGE, width=2.0)
    add_arrow(slide, 3.02, 4.40, 3.02, 3.40, color=ORANGE, width=2.0)
    add_text(slide, "Hướng tiếp tuyến đổi tức thời tại V", 1.24, 1.80, 4.30, 0.48, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "κ(t) = [x′y″ − y′x″] / (x′² + y′²)³ᐟ²", 6.48, 1.65, 5.90, 0.86, size=20)
    add_text(slide, "Tại đỉnh polyline, đạo hàm theo hướng không liên tục ⇒ độ cong κ rất lớn hoặc không xác định.", 6.75, 2.83, 5.25, 0.95, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 6.62, 4.10, 1.72, 1.58, title="Giảm tốc", body="Bộ điều khiển phải hãm mạnh.", accent=ORANGE, title_size=14, body_size=11.5)
    add_card(slide, 8.58, 4.10, 1.72, 1.58, title="Quay tại chỗ", body="Dừng–quay–đi.", accent=RED, title_size=14, body_size=11.5)
    add_card(slide, 10.54, 4.10, 1.72, 1.58, title="Lệch bám", body="Khó theo đúng góc gãy.", accent=PURPLE, title_size=14, body_size=11.5)
    add_text(slide, "Bài toán không chỉ là “đường nhìn đẹp hơn”, mà là tạo chuyển hướng có thể chạy.", 6.58, 6.10, 5.68, 0.42, size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # 08 — Baselines
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Ba bộ làm mượt Nav2 dùng làm đối chứng", "Hệ thống", 8, subtitle="Cùng nhận đúng một đường do bộ lập kế hoạch tạo trong mỗi nhóm ghép cặp")
    cols = [
        ("SIMPLE", BLUE, "Lặp cục bộ trên điểm đường", "yᵢ ← yᵢ + w_d(xᵢ−yᵢ)\n+ w_s(yᵢ₋₁+yᵢ₊₁−2yᵢ)", "Nhanh; giảm góc gãy nhỏ.\nKhông dùng hình bao robot trong phép cập nhật."),
        ("SAVITZKY–GOLAY", ORANGE, "Lọc đa thức cửa sổ", "ŷᵢ = (−2yᵢ₋₃+3yᵢ₋₂+6yᵢ₋₁\n+7yᵢ+6yᵢ₊₁+3yᵢ₊₂−2yᵢ₊₃)/21", "Lọc dao động cao tần.\nKhó xử lý góc gãy lớn."),
        ("CONSTRAINED", PURPLE, "Tối ưu có bản đồ chi phí", "J = w_sJ_s + w_κJ_κ\n+ w_dJ_d + w_cJ_cost", "Tối ưu nặng hơn. Cấu hình thử:\nw_s=200000; w_cost=0,015; w_κ=w_d=0."),
    ]
    for i, (name, accent, subtitle, formula, note) in enumerate(cols):
        x = 0.62 + i * 4.23
        add_rect(slide, x, 1.52, 3.86, 4.95, fill=WHITE, line=LINE)
        add_rect(slide, x, 1.52, 3.86, 0.52, fill=accent, line=None, radius=False)
        add_text(slide, name, x, 1.55, 3.86, 0.40, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, subtitle, x + 0.25, 2.25, 3.36, 0.40, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_equation(slide, formula, x + 0.28, 2.93, 3.30, 1.32, color=NAVY, fill=LIGHT, size=13.5 if i != 1 else 11.5)
        add_text(slide, note, x + 0.30, 4.65, 3.25, 1.15, size=12.2, color=MUTED, align=PP_ALIGN.CENTER)
    add_pill(slide, "KHOẢNG TRỐNG", 5.14, 6.67, 2.98, fill=PALE_RED, color=RED, size=10.5)
    add_text(slide, "Các phương pháp đối chứng không đồng thời mô hình hóa G², bánh xe vi sai, điều kiện ưu thế thời gian và quy hoạch động.", 1.70, 6.92, 9.94, 0.30, size=11.9, color=MUTED, align=PP_ALIGN.CENTER)

    # 09 — Goal constraints
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Mục tiêu thiết kế của PSTMO", "Hệ thống", 9)
    add_equation(slide, "Tìm chuỗi trạng thái z₁…z_N có chi phí nhỏ nhất\ntrên tập phương án đã đạt toàn bộ điều kiện bắt buộc", 1.55, 1.55, 10.25, 1.05, color=NAVY, fill=LIGHT, size=20)
    constraints = [
        ("G²", "κ liên tục tại mối nối", BLUE),
        ("v_L, v_R", "Không đảo bánh trong đoạn chuyển tiếp", PURPLE),
        ("F_map", "Vùng quét hình bao không va chạm", GREEN),
        ("T", "Đoạn chuyển tiếp có lợi hơn quay tại chỗ", ORANGE),
        ("d_i+d_{i+1}", "Vùng cắt không chồng lấn", TEAL),
    ]
    for i, (symbol, text, accent) in enumerate(constraints):
        x = 0.62 + i * 2.52
        add_rect(slide, x, 3.20, 2.18, 2.15, fill=WHITE, line=LINE)
        add_text(slide, symbol, x + 0.12, 3.47, 1.94, 0.52, size=23, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, text, x + 0.20, 4.20, 1.78, 0.70, size=12.0, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Phương án xử lý tại góc = một cách xử lý cụ thể, xác định bởi trạng thái, khoảng cắt d và tỷ lệ α=q/d.", 1.15, 5.98, 11.05, 0.64, size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # Section II
    add_section_slide(prs, "II", "PHƯƠNG PHÁP PSTMO", "Từ điều kiện hóa đường tới ghép đường cuối cùng — mỗi quyết định đều có công thức và điều kiện kiểm tra.", TEAL, VISUALS_3D / "swept_footprint_clean.png")

    # 11 — Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Một quy trình duy nhất, theo đúng thứ tự kiểm tra", "PSTMO", 11)
    add_picture_contain(slide, ASSETS / "figures" / "figure_01_pipeline.png", 0.70, 1.48, 11.95, 2.95)
    steps = [
        ("01", "Điều kiện hóa", "Bỏ điểm trùng, RDP an toàn, triệt zíc-zắc."),
        ("02", "Tạo phương án", "Quay tại chỗ hoặc Bézier bậc năm G²."),
        ("03", "Điều kiện bắt buộc", "Hình học, bánh xe, hình bao, thời gian."),
        ("04", "Chọn toàn cục", "Ưu thế thời gian, chi phí và quy hoạch động."),
    ]
    for i, (n, title, body) in enumerate(steps):
        x = 0.75 + i * 3.12
        add_circle_label(slide, n, x, 4.82, 0.48, fill=[BLUE, TEAL, GREEN, PURPLE][i], size=10.5)
        add_text(slide, title, x + 0.62, 4.80, 2.22, 0.33, size=14.5, color=NAVY, bold=True)
        add_text(slide, body, x + 0.02, 5.35, 2.82, 0.78, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_equation(slide, "Nếu kiểm tra cuối thất bại → FailedToSmoothPath; không âm thầm trả lại đường gốc.", 1.72, 6.38, 9.90, 0.54, color=RED, fill=PALE_RED, size=14)

    # 12 — Conditioning
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Tiền xử lý: giảm nhiễu nhưng không phá hành lang", "PSTMO", 12, subtitle="Điều kiện hóa đường (path conditioning) · chế độ hiện tại: condition_only")
    add_picture_contain(slide, ASSETS / "figures" / "figure_03_conditioning_actual.png", 0.72, 1.50, 7.32, 4.75)
    add_card(slide, 8.38, 1.56, 4.20, 1.34, title="1. Bỏ điểm trùng", body="Tránh đoạn có chiều dài bằng 0 và đạo hàm suy biến.", accent=BLUE)
    add_card(slide, 8.38, 3.08, 4.20, 1.50, title="2. RDP có điều kiện an toàn", body="Chỉ nhận dây cung khi sai số ≤ ε_RDP và vùng quét hình bao an toàn.", accent=TEAL)
    add_card(slide, 8.38, 4.78, 4.20, 1.47, title="3. Triệt zíc-zắc cục bộ", body="Loại chuỗi đổi dấu góc do lưới, nhưng vẫn giữ điểm đầu/đích.", accent=PURPLE)
    add_equation(slide, "ε_RDP = 1,5 × độ phân giải = 1,5 × 0,05 = 0,075 m", 1.37, 6.46, 10.60, 0.54, color=TEAL, fill=PALE_GREEN, size=15)

    # 13 — Corner states
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Phát hiện góc và ba trạng thái xử lý", "PSTMO", 13)
    add_equation(slide, "θ = atan2(u × v, u · v)", 0.82, 1.52, 4.42, 0.78, size=22)
    add_text(slide, "u: hướng cạnh vào · v: hướng cạnh ra · dấu θ phân biệt quay trái/phải", 0.95, 2.47, 4.18, 0.72, size=12.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_equation(slide, "|θ| < 5°  →  giữ nguyên góc\n5° ≤ |θ| < 170°  →  xét chuyển tiếp hoặc quay tại chỗ", 0.82, 3.50, 4.42, 1.30, color=NAVY, fill=PALE_BLUE, size=16)
    states = [
        ("PASS-THROUGH", "Giữ nguyên góc", "Góc nhỏ; không cần đoạn chuyển tiếp.", BLUE),
        ("TRANSITION", "Đoạn chuyển tiếp G²", "Cắt hai cạnh và chèn Bézier.", TEAL),
        ("PIVOT", "Quay tại chỗ", "Dừng tại đỉnh, đổi góc hướng rồi đi tiếp.", ORANGE),
    ]
    for i, (eng, vi, body, accent) in enumerate(states):
        x = 5.65 + i * 2.40
        add_rect(slide, x, 1.55, 2.14, 4.58, fill=WHITE, line=LINE)
        add_pill(slide, eng, x + 0.20, 1.82, 1.74, fill=accent, color=WHITE, size=9)
        if i == 0:
            add_line(slide, x + 0.45, 3.22, x + 1.70, 3.22, color=accent, width=5)
        elif i == 1:
            add_line(slide, x + 0.38, 3.55, x + 0.92, 3.55, color=accent, width=4)
            arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x + 0.83), Inches(2.75), Inches(0.80), Inches(0.82))
            arc.line.color.rgb = rgb(accent); arc.line.width = Pt(4); arc.fill.background()
            add_line(slide, x + 1.62, 2.78, x + 1.62, 2.32, color=accent, width=4)
        else:
            add_line(slide, x + 0.38, 3.55, x + 1.07, 3.55, color=accent, width=4)
            add_circle_label(slide, "↻", x + 0.88, 2.72, 0.50, fill=accent, size=16)
            add_line(slide, x + 1.25, 3.37, x + 1.74, 2.90, color=accent, width=4)
        add_text(slide, vi, x + 0.20, 4.13, 1.74, 0.52, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.20, 4.82, 1.74, 0.70, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "Ở mỗi góc, PSTMO tạo các phương án cụ thể rồi kiểm tra, so sánh và lựa chọn.", 1.18, 6.53, 11.00, 0.36, size=13.2, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # 14 — G1 / G2 and quintic rationale
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "G¹, G² và lý do chọn Bézier bậc năm", "PSTMO", 14)
    add_rect(slide, 0.72, 1.52, 3.72, 2.16, fill=PALE_ORANGE, line=LINE)
    add_text(slide, "G¹ · LIÊN TỤC TIẾP TUYẾN", 0.98, 1.79, 3.20, 0.36, size=15.5, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Cùng hướng tại điểm nối\nnhưng κ có thể nhảy.", 1.12, 2.40, 2.92, 0.74, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.76, 1.52, 3.72, 2.16, fill=PALE_GREEN, line=LINE)
    add_text(slide, "G² · LIÊN TỤC ĐỘ CONG", 5.02, 1.79, 3.20, 0.36, size=15.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Cùng hướng và cùng κ.\nĐoạn thẳng có κ=0.", 5.16, 2.40, 2.92, 0.74, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.80, 1.52, 3.82, 2.16, fill=PALE_BLUE, line=LINE)
    add_text(slide, "YÊU CẦU NỐI", 9.08, 1.79, 3.25, 0.36, size=15.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "κ(0)=κ(1)=0", 9.24, 2.32, 2.90, 0.66, color=NAVY, fill=WHITE, size=22)
    table_data = [
        ["Bậc", "Điểm điều khiển", "Khớp hướng", "κ=0 hai đầu", "Còn tham số hình dạng"],
        ["2", "3", "Có", "Không", "Không"],
        ["3–4", "4–5", "Có", "Có thể", "Rất hạn chế"],
        ["5", "6", "Có", "Có theo cấu trúc", "Có: α=q/d"],
    ]
    add_table(slide, table_data, 0.88, 4.08, 11.58, 1.92, col_widths=[1.0, 2.05, 1.75, 2.28, 4.50], font_size=11.5)
    add_text(slide, "Bậc năm là cấu trúc gọn nhất trong thiết kế này để vừa áp điều kiện biên G², vừa giữ một tham số tối ưu hình dạng.", 1.13, 6.43, 11.08, 0.50, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 15 — Bezier construction
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Xây dựng đoạn chuyển tiếp Bézier bậc năm", "PSTMO", 15)
    add_picture_contain(slide, ASSETS / "figures" / "figure_04_bezier_g2.png", 0.72, 1.48, 5.80, 4.80)
    add_equation(slide, "B(t)=Σᵢ₌₀⁵ C(5,i)(1−t)⁵⁻ⁱtⁱPᵢ,   0≤t≤1", 6.88, 1.52, 5.72, 0.78, size=18)
    add_equation(slide, "A=V−d·u     B=V+d·v     q=αd", 6.88, 2.55, 5.72, 0.70, color=TEAL, fill=PALE_GREEN, size=19)
    add_equation(slide, "P₀=A      P₁=A+q·u      P₂=A+2q·u\nP₃=B−2q·v   P₄=B−q·v      P₅=B", 6.88, 3.50, 5.72, 1.18, color=NAVY, fill=LIGHT, size=17)
    add_card(slide, 6.92, 4.98, 2.58, 1.30, title="d · khoảng cắt", body="Quy mô vùng chuyển hướng.", accent=BLUE, title_size=13.5, body_size=11.2)
    add_card(slide, 9.83, 4.98, 2.58, 1.30, title="α=q/d · tỷ lệ hình dạng", body="Phân bố độ cong bên trong.", accent=PURPLE, title_size=13.5, body_size=11.2)
    add_caption(slide, "Hình sinh từ đúng cấu trúc Bézier G² trong mã nguồn", 0.86, 6.36, 5.52)

    # 16 — G2 proof
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Chứng minh mối nối đạt G²", "PSTMO", 16, subtitle="Ba bước ngắn: vị trí → tiếp tuyến → độ cong")
    proof = [
        ("1", "Khớp vị trí", "B(0)=P₀=A\nB(1)=P₅=B", BLUE),
        ("2", "Khớp tiếp tuyến", "B′(0)=5q·u\nB′(1)=5q·v", TEAL),
        ("3", "Triệt đạo hàm bậc hai", "B″(0)=20(P₂−2P₁+P₀)=0\nB″(1)=20(P₅−2P₄+P₃)=0", PURPLE),
    ]
    for i, (n, title, formula, accent) in enumerate(proof):
        x = 0.72 + i * 4.18
        add_rect(slide, x, 1.62, 3.82, 3.05, fill=WHITE, line=LINE)
        add_circle_label(slide, n, x + 0.20, 1.84, 0.48, fill=accent)
        add_text(slide, title, x + 0.82, 1.85, 2.68, 0.38, size=15.5, color=NAVY, bold=True)
        add_equation(slide, formula, x + 0.24, 2.60, 3.34, 1.42, color=accent, fill=LIGHT, size=15 if i != 2 else 12.8)
    add_equation(slide, "κ = (B′ × B″) / ‖B′‖³   ⇒   κ(0)=κ(1)=0", 2.05, 5.12, 9.22, 0.82, color=TEAL, fill=PALE_GREEN, size=22)
    add_text(slide, "Vì đoạn thẳng kề cũng có κ=0, độ cong không nhảy tại hai điểm nối.", 2.04, 6.22, 9.25, 0.42, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Lưu ý: G² là độ mượt hình học; không tự động bảo đảm jerk theo thời gian bằng 0.", 2.03, 6.70, 9.28, 0.32, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 17 — d and alpha
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Khoảng cắt d và tỷ lệ α có hai vai trò khác nhau", "PSTMO", 17)
    add_rect(slide, 0.72, 1.55, 5.85, 4.90, fill=PALE_BLUE, line=LINE)
    add_text(slide, "d · TRIM DISTANCE · KHOẢNG CẮT", 1.03, 1.86, 5.22, 0.42, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Quyết định đoạn chuyển tiếp chiếm bao nhiêu chiều dài trên hai cạnh.", 1.30, 2.48, 4.70, 0.58, size=15, color=INK, align=PP_ALIGN.CENTER)
    add_equation(slide, "d_pref = min(d_max, L_in, L_out)", 1.25, 3.32, 4.80, 0.70, color=NAVY, fill=WHITE, size=18)
    add_equation(slide, "d_compat = min(d_pref, b_in, b_out)", 1.25, 4.26, 4.80, 0.70, color=NAVY, fill=WHITE, size=18)
    add_text(slide, "Tối đa hai giá trị: ưu tiên và tương thích với đoạn chung.", 1.40, 5.40, 4.50, 0.46, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.78, 1.55, 5.85, 4.90, fill=PALE_PURPLE, line=LINE)
    add_text(slide, "α=q/d · SHAPE RATIO · TỶ LỆ HÌNH DẠNG", 7.08, 1.86, 5.25, 0.42, size=17, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Khi d cố định, α điều chỉnh đa giác điều khiển và phân bố độ cong.", 7.35, 2.48, 4.70, 0.58, size=15.5, color=INK, align=PP_ALIGN.CENTER)
    add_equation(slide, "q = αd,     0 < α ≤ 0,5", 7.32, 3.32, 4.80, 0.70, color=NAVY, fill=WHITE, size=20)
    add_equation(slide, "E_κ = ∫₀ᴸ κ(s)² ds", 7.32, 4.26, 4.80, 0.70, color=PURPLE, fill=WHITE, size=20)
    add_text(slide, "Trong cùng một d, chọn α khả thi có Eκ nhỏ nhất.", 7.48, 5.40, 4.50, 0.46, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "Eκ là tích phân bình phương độ cong, tức chỉ số uốn hình học; không phải điện năng đo từ pin.", 1.55, 6.70, 10.25, 0.32, size=11.8, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 18 — alpha search
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Tìm α theo lưới phân cấp", "PSTMO", 18, subtitle="Tìm thô → phục hồi khi cần → tinh chỉnh quanh phương án có Eκ nhỏ nhất")
    add_picture_contain(slide, ASSETS / "figures" / "figure_05_alpha_search.png", 0.72, 1.48, 7.62, 4.98)
    stages = [
        ("COARSE · THÔ", "α={0,1; 0,2; 0,3; 0,4; 0,5}", BLUE),
        ("RECOVERY · PHỤC HỒI", "Chỉ chạy nếu toàn bộ lưới thô thất bại:\n{0,15; 0,25; 0,35; 0,45}", RED),
        ("REFINEMENT · TINH", "Chia khoảng quanh phương án có Eκ nhỏ nhất thành 10 khoảng, 11 nút.", GREEN),
    ]
    for i, (title, body, accent) in enumerate(stages):
        add_card(slide, 8.64, 1.58 + i * 1.56, 3.92, 1.32, title=title, body=body, accent=accent, title_size=12.5, body_size=11.2)
    add_equation(slide, "Chỉ so Eκ giữa các α đã đạt mọi điều kiện bắt buộc.", 8.64, 6.25, 3.92, 0.58, color=TEAL, fill=PALE_GREEN, size=12.8)

    # 19 — Hard gates
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Kiểm tra tính khả thi của từng phương án", "PSTMO", 19, subtitle="Chỉ các phương án đạt toàn bộ điều kiện bắt buộc mới được so sánh")
    gates = [
        ("1 · HÌNH HỌC", "B′ không suy biến; κ hữu hạn; không xuất hiện đổi dấu quay ngoài ý muốn.", BLUE),
        ("2 · ĐỘNG HỌC", "Bánh trong không chạy lùi; v, ω và gia tốc đều nằm trong giới hạn.", PURPLE),
        ("3 · AN TOÀN", "Chi phí tại tâm đạt ngưỡng; vùng quét hình bao không giao vật cản.", GREEN),
        ("4 · THỜI GIAN", "Biểu đồ vận tốc hội tụ; đoạn chuyển tiếp nhanh hơn quay tại chỗ theo biên 0,15 s.", ORANGE),
        ("5 · HẬU KIỂM", "Đường cuối giữ đầu–đích, không chồng vùng cắt và vẫn an toàn sau khi ghép.", TEAL),
    ]
    for i, (title, body, accent) in enumerate(gates):
        x = 0.50 + i * 2.55
        add_card(slide, x, 1.55, 2.28, 4.36, title=title, body=body, accent=accent, title_size=11.5, body_size=11.0)
        if i < len(gates) - 1:
            add_text(slide, "→", x + 2.28, 3.22, 0.27, 0.38, size=19, color=LINE, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "Đạt cả 5 nhóm điều kiện  →  mới được tính J và so sánh phương án", 1.24, 6.10, 10.85, 0.52, color=TEAL, fill=PALE_GREEN, size=14.5)
    add_text(slide, "Va chạm là điều kiện loại bắt buộc; Eκ nhỏ không thể bù cho một phương án không an toàn.", 1.05, 6.76, 11.05, 0.30, size=11.6, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 20 — Differential drive
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Cổng động học robot vi sai", "PSTMO", 20)
    add_picture_contain(slide, VISUALS_3D / "wheel_layout_clean.png", 0.70, 1.50, 5.35, 4.16)
    add_pill(slide, "MÔ HÌNH TỪ STL/SDF", 1.72, 5.82, 3.30, fill=PALE_BLUE, color=BLUE, size=10)
    add_text(slide, "Bánh trái: y=+0,1274 m · bánh phải: y=−0,1274 m · b=0,2548 m", 0.90, 6.26, 4.95, 0.30, size=10.8, color=MUTED, align=PP_ALIGN.CENTER)
    add_equation(slide, "v_L = v(1 − bκ/2)\nv_R = v(1 + bκ/2)\nω = vκ", 6.48, 1.58, 5.82, 1.66, size=22)
    add_equation(slide, "1 − |bκ|/2 ≥ 0   ⇒   |κ| ≤ 2/b", 6.48, 3.62, 5.82, 0.76, color=PURPLE, fill=PALE_PURPLE, size=20)
    add_text(slide, "Đoạn chuyển tiếp là pha tịnh tiến liên tục nên bánh trong không được chạy lùi. Quay tại chỗ là trạng thái riêng, cho phép hai bánh quay ngược chiều.", 6.78, 4.78, 5.23, 1.02, size=14, color=INK, align=PP_ALIGN.CENTER)
    add_pill(slide, "GIỚI HẠN THỬ NGHIỆM", 7.75, 6.04, 3.34, fill=PALE_ORANGE, color=ORANGE, size=10)
    add_text(slide, "b mô hình=0,2548 m · b hiệu dụng Gazebo=0,2834 m · |v_bánh|≤0,36 m/s", 6.56, 6.47, 5.70, 0.34, size=10.8, color=MUTED, align=PP_ALIGN.CENTER)

    # 21 — Footprint
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Điều kiện an toàn: hình bao robot và vùng quét", "PSTMO", 21)
    add_picture_contain(slide, VISUALS_3D / "swept_footprint_clean.png", 0.72, 1.50, 5.35, 4.98)
    add_caption(slide, "Màu xanh lá: hợp các hình bao · màu xanh đậm: vật cản", 0.92, 6.52, 4.96)
    add_equation(slide, "F_map(x,y,ψ) = { [x,y]ᵀ + R(ψ)f  |  f∈F_body }", 6.46, 1.58, 5.90, 0.80, size=17.5)
    add_card(slide, 6.50, 2.72, 5.80, 1.10, title="Tư thế an toàn", body="Đa giác hình bao tại một tư thế không chạm ô gây va chạm, ô chưa biết hoặc ngoài bản đồ.", accent=BLUE, title_size=14, body_size=10.9)
    add_card(slide, 6.50, 4.04, 5.80, 1.10, title="Vùng quét an toàn", body="Nội suy vị trí và góc hướng giữa các tư thế; kiểm tra toàn bộ đa giác hình bao.", accent=GREEN, title_size=14, body_size=10.9)
    add_equation(slide, "Chi phí ô tại tâm ≤252  ∧  hình bao đặc không va chạm", 6.50, 5.55, 5.80, 0.68, color=TEAL, fill=PALE_GREEN, size=15)
    add_text(slide, "Hình bao thử nghiệm: 0,44 × 0,34 m · độ phân giải bản đồ chi phí: 0,05 m", 6.54, 6.54, 5.70, 0.34, size=10.8, color=MUTED, align=PP_ALIGN.CENTER)

    # 22 — Curvature-dependent speed limit
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Từ độ cong tới giới hạn vận tốc", "PSTMO", 22, subtitle="Giới hạn vận tốc phụ thuộc độ cong và biểu đồ vận tốc dọc đường")
    add_equation(slide, "v_limit(κ)=min{ v_max,  ω_max/|κ|,  √(a_y,max/|κ|),\n v_w,max / max(|1−bκ/2|, |1+bκ/2|) }", 0.75, 1.52, 6.05, 1.30, size=16.5)
    add_text(slide, "Độ cong càng lớn → giới hạn vận tốc tịnh tiến càng thấp.", 1.00, 3.10, 5.55, 0.48, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "Quét tiến:  vᵢ ≤ √(vᵢ₋₁² + 2a_accΔs)\nQuét lùi:   vᵢ₋₁ ≤ √(vᵢ² + 2a_decΔs)", 0.75, 3.85, 6.05, 1.20, color=NAVY, fill=LIGHT, size=16.5)
    add_text(slide, "Biểu đồ vận tốc là dãy v(s) sau khi áp giới hạn vận tốc cục bộ và các giới hạn tăng tốc, giảm tốc.", 0.98, 5.43, 5.58, 0.88, size=13.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_picture_contain(slide, ASSETS / "figures" / "figure_07_kinematic_time_gate.png", 7.12, 1.50, 5.48, 4.85)
    add_caption(slide, "Ví dụ biểu đồ vận tốc tịnh tiến và vận tốc góc dùng để tính thời gian", 7.30, 6.45, 5.12)

    # 23 — Time gate
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "So sánh đoạn chuyển tiếp với quay tại chỗ", "PSTMO", 23, subtitle="Điều kiện ưu thế thời gian (time gate)")
    add_rect(slide, 0.72, 1.48, 5.50, 4.88, fill=PALE_ORANGE, line=LINE)
    add_text(slide, "PIVOT · QUAY TẠI CHỖ", 1.05, 1.82, 4.84, 0.38, size=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Đi tới đỉnh → giảm về 0 → quay → tăng tốc rời đỉnh", 1.10, 2.45, 4.75, 0.62, size=15, color=INK, align=PP_ALIGN.CENTER)
    add_equation(slide, "T_rot = 2√(|θ|/a_ω)  nếu chưa đạt ω_max", 1.20, 3.32, 4.50, 0.72, color=ORANGE, fill=WHITE, size=16)
    add_equation(slide, "T_pivot = T_đến + T_rot + T_rời", 1.20, 4.28, 4.50, 0.72, color=NAVY, fill=WHITE, size=17)
    add_rect(slide, 6.54, 1.48, 6.05, 4.88, fill=PALE_GREEN, line=LINE)
    add_text(slide, "TRANSITION · ĐOẠN CHUYỂN TIẾP", 6.90, 1.82, 5.35, 0.38, size=17, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "Δtᵢ = 2Δsᵢ/(vᵢ+vᵢ₊₁)\nT_curve = Σᵢ Δtᵢ", 7.35, 2.50, 4.45, 1.02, color=TEAL, fill=WHITE, size=19)
    add_equation(slide, "T_fastest + ΔT < T_pivot\nΔT = 0,15 s", 7.20, 3.88, 4.75, 1.08, color=GREEN, fill=WHITE, size=21)
    add_text(slide, "Nếu quay tại chỗ không an toàn, đoạn chuyển tiếp vẫn được giữ khi bản thân nó khả thi.", 7.15, 5.30, 4.85, 0.54, size=12.6, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "Tất cả phương án được đặt trên cùng một cửa sổ hình học trước khi so sánh thời gian.", 1.50, 6.66, 10.35, 0.34, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 24 — Objective
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Hàm mục tiêu dùng để lựa chọn phương án", "PSTMO", 24, subtitle="Chỉ so sánh sau khi phương án đạt điều kiện bắt buộc và điều kiện ưu thế thời gian")
    add_equation(slide, "r_cost = min(1, peak_cost/252)\nr_ω = min(1, |ω|_max/ω_max)\neκ = (Eκ/E_ref)/(Eκ/E_ref+1)\nE_ref = 1 m⁻¹", 0.74, 1.43, 5.20, 2.12, size=15.8)
    add_equation(slide, "J = 0,15·r_cost + 0,10·r_ω + 0,75·eκ  →  min", 0.74, 3.78, 5.20, 0.80, color=PURPLE, fill=PALE_PURPLE, size=16.5)
    add_text(slide, "J nhỏ hơn → phương án tại góc được ưu tiên trong quy hoạch động.", 0.98, 4.78, 4.75, 0.64, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    labels = [("CHI PHÍ BẢN ĐỒ", 0.15, RED), ("VẬN TỐC GÓC", 0.10, ORANGE), ("CHỈ SỐ UỐN Eκ", 0.75, PURPLE)]
    y = 1.72
    for name, weight, accent in labels:
        add_text(slide, name, 6.52, y, 2.05, 0.32, size=13, color=NAVY, bold=True)
        add_rect(slide, 8.56, y + 0.02, 3.32, 0.28, fill=LIGHT, line=None, radius=False)
        add_rect(slide, 8.56, y + 0.02, 3.32 * weight / 0.75, 0.28, fill=accent, line=None, radius=False)
        add_text(slide, f"{weight:.2f}".replace(".", ","), 11.98, y - 0.02, 0.52, 0.34, size=13, color=accent, bold=True, align=PP_ALIGN.RIGHT)
        y += 1.12
    add_card(slide, 6.48, 5.22, 5.90, 1.12, title="Không nhầm chi phí mềm với an toàn", body="peak_cost chỉ dùng để so sánh phương án đã an toàn; va chạm luôn bị loại trước.", accent=RED, title_size=13.5, body_size=11.5)
    add_text(slide, "Cửa sổ cạnh tranh tối đa: 10 s · Eκ có đơn vị m⁻¹ nhưng eκ và J không có đơn vị.", 1.58, 6.68, 10.15, 0.32, size=11.6, color=MUTED, align=PP_ALIGN.CENTER)

    # 25 — DP
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Quy hoạch động chọn chuỗi trạng thái toàn đường", "PSTMO", 25)
    add_picture_contain(slide, ASSETS / "figures" / "figure_06_two_trim_dp.png", 0.72, 1.48, 7.10, 4.92)
    add_equation(slide, "d(zᵢ)+d(zᵢ₊₁)+m ≤ Lᵢ\nm = 0,05 m", 8.18, 1.62, 4.35, 0.90, color=TEAL, fill=PALE_GREEN, size=17)
    add_text(slide, "Điều kiện tương thích:\nhai vùng cắt không chồng lấn trên đoạn chung.", 8.40, 2.60, 3.90, 0.78, size=14.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_equation(slide, "Dᵢ(z)=Jᵢ(z)+min_{z′ tương thích} Dᵢ₋₁(z′)", 8.18, 3.72, 4.35, 0.98, color=PURPLE, fill=PALE_PURPLE, size=16)
    add_card(slide, 8.18, 5.03, 4.35, 1.20, title="Độ phức tạp", body="O(NK²), với N góc và tối đa K trạng thái mỗi góc.", accent=BLUE, title_size=13.5, body_size=11.5)
    add_text(slide, "Phương án có J nhỏ tại một góc chưa chắc tạo thành đường tốt: quy hoạch động giải xung đột giữa các góc.", 1.10, 6.68, 11.10, 0.36, size=12.6, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 26 — Output/controller
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Ghép đường cuối và bàn giao cho bộ điều khiển", "PSTMO", 26)
    checks = [
        ("Giữ điểm đầu–đích", "Vị trí đầu, vị trí đích và hướng đích không đổi.", BLUE),
        ("Nội suy đường thẳng", "Khoảng lấy mẫu đầu ra: 0,05 m.", TEAL),
        ("Chèn đoạn chuyển tiếp", "Góc hướng lấy theo tiếp tuyến Bézier.", PURPLE),
        ("Biểu diễn quay tại chỗ", "Hai tư thế cùng vị trí nhưng khác góc hướng.", ORANGE),
        ("Hậu kiểm toàn đường", "Đầu–đích, thời gian, chồng lấn và vùng quét.", GREEN),
    ]
    for i, (title, body, accent) in enumerate(checks):
        y = 1.50 + i * 0.96
        add_circle_label(slide, "✓", 0.86, y + 0.04, 0.42, fill=accent, size=12)
        add_text(slide, title, 1.50, y, 2.52, 0.32, size=14.5, color=NAVY, bold=True)
        add_text(slide, body, 4.10, y, 4.07, 0.50, size=12.5, color=MUTED)
    add_rect(slide, 8.62, 1.52, 3.80, 4.80, fill=LIGHT, line=LINE)
    add_node(slide, "PSTMO\nnav_msgs/Path", 9.42, 2.08, 2.20, 0.82, fill=TEAL, size=13)
    add_arrow(slide, 10.52, 3.05, 10.52, 3.72, color=LINE, width=2.2)
    add_node(slide, "Bộ điều khiển RPP\n20 Hz", 9.42, 3.78, 2.20, 0.82, fill=ORANGE, size=12)
    add_arrow(slide, 10.52, 4.75, 10.52, 5.42, color=LINE, width=2.2)
    add_node(slide, "cmd_vel\n(v, ω)", 9.42, 5.48, 2.20, 0.58, fill=GREEN, size=12)
    add_equation(slide, "Bộ làm mượt chỉnh đường · Bộ điều khiển phát lệnh", 2.08, 6.54, 9.10, 0.54, color=NAVY, fill=PALE_BLUE, size=15)

    # Section III
    add_section_slide(prs, "III", "THỰC NGHIỆM & KẾT QUẢ", "Hình 3D dùng để giải thích mô hình; kết quả được chứng minh bằng ảnh Gazebo, RViz2 và dữ liệu ROS đã lưu.", GREEN, VISUALS_3D / "robot_gazebo_warehouse_hero.png")

    # 28 — System model
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Mô hình mô phỏng gồm những gì?", "Thực nghiệm", 28)
    add_picture_contain(slide, VISUALS_3D / "robot_isometric_clean.png", 0.62, 1.46, 3.92, 2.90)
    add_picture_contain(slide, ROOT / "results" / "gui_validation_20260724" / "gazebo_research_warehouse.png", 4.70, 1.46, 3.92, 2.90)
    add_picture_contain(slide, ROOT / "results" / "gui_validation_20260724" / "rviz_research_ui_final.png", 8.78, 1.46, 3.92, 2.90)
    add_caption(slide, "3D từ STL/SDF: cấu tạo", 0.78, 4.43, 3.60)
    add_caption(slide, "Gazebo: trạng thái vật lý và môi trường", 4.86, 4.43, 3.60)
    add_caption(slide, "RViz2: dữ liệu điều hướng và đường đi", 8.94, 4.43, 3.60)
    components = [
        ("Robot", "Vi sai; hình bao 0,44×0,34 m", BLUE),
        ("Cảm biến & định vị", "LiDAR 2D + AMCL", PURPLE),
        ("Điều hướng", "5 bộ lập kế hoạch + 5 phương án đường", TEAL),
        ("Bám đường & an toàn", "RPP + làm mượt vận tốc + giám sát va chạm", GREEN),
    ]
    for i, (title, body, accent) in enumerate(components):
        add_card(slide, 0.72 + i * 3.12, 5.05, 2.83, 1.42, title=title, body=body, accent=accent, title_size=12.4, body_size=10.2)
    add_text(slide, "Chỉ ảnh Gazebo/RViz2 và dữ liệu ROS được dùng làm bằng chứng thực nghiệm; hình 3D dùng để giải thích mô hình.", 1.23, 6.70, 10.90, 0.30, size=10.8, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    # 29 — Environments
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Bảy môi trường thử nghiệm trong Gazebo/RViz2", "Thực nghiệm", 29)
    envs = [
        ("open_arena.png", "Không gian mở"),
        ("research_warehouse.png", "Kho nghiên cứu"),
        ("narrow_aisles.png", "Lối đi hẹp"),
        ("office_maze.png", "Mê cung văn phòng"),
        ("warehouse_cross_aisles.png", "Kho giao cắt"),
        ("warehouse_dispatch.png", "Kho điều phối"),
        ("warehouse_long_aisles.png", "Kho lối đi dài"),
    ]
    positions = [(0.55, 1.45), (3.63, 1.45), (6.71, 1.45), (9.79, 1.45), (2.09, 4.03), (5.17, 4.03), (8.25, 4.03)]
    for (file, label), (x, y) in zip(envs, positions):
        add_picture_contain(slide, ASSETS / "gazebo" / file, x, y, 2.75, 1.90)
        add_text(slide, label, x, y + 1.97, 2.75, 0.28, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Mỗi môi trường dùng một cặp đầu–đích đại diện; cùng một đường gốc được chia sẻ giữa năm phương án của mỗi bộ lập kế hoạch.", 1.22, 6.72, 10.90, 0.34, size=12.1, color=MUTED, align=PP_ALIGN.CENTER)

    # 30 — Protocol
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Thiết kế thực nghiệm và cách ghép cặp", "Thực nghiệm", 30)
    add_picture_contain(slide, ASSETS / "figures" / "figure_12_test_matrix.png", 0.70, 1.46, 7.38, 4.86)
    add_metric(slide, 8.42, 1.56, 1.82, 1.48, value="7", label="môi trường", accent=BLUE)
    add_metric(slide, 10.52, 1.56, 1.82, 1.48, value="5", label="bộ lập kế hoạch", accent=PURPLE)
    add_metric(slide, 8.42, 3.30, 1.82, 1.48, value="5", label="phương án", accent=TEAL)
    add_metric(slide, 10.52, 3.30, 1.82, 1.48, value="175", label="lượt thực thi", accent=GREEN)
    add_card(slide, 8.42, 5.12, 3.92, 1.18, title="Đơn vị ghép cặp", body="Môi trường + bộ lập kế hoạch + SHA-256 của đúng đường gốc.", accent=ORANGE, title_size=13, body_size=10.8)
    add_text(slide, "Mỗi lượt chạy trong một mô phỏng mới; số lần lặp hiện tại: 1 cho mỗi tổ hợp.", 1.18, 6.66, 11.00, 0.34, size=12.5, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 31 — Evidence
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "RViz2 và Gazebo trả về dữ liệu gì?", "Thực nghiệm", 31)
    rviz = ASSETS / "rviz_cases" / "research_warehouse__lower_left_diagonal__Smac2D.png"
    case = ASSETS / "execution_cases" / "research_warehouse_smac2d.png"
    add_picture_contain(slide, rviz, 0.68, 1.46, 6.15, 4.95)
    add_picture_contain(slide, case, 7.06, 1.46, 5.60, 4.95)
    add_caption(slide, "RViz2: đường gốc/PSTMO, hình bao, bản đồ chi phí và dữ liệu chẩn đoán", 0.86, 6.47, 5.80)
    add_caption(slide, "Đối chiếu kế hoạch với quỹ đạo thực thi trong Gazebo cho đủ năm phương án", 7.23, 6.47, 5.25)
    add_text(slide, "Mỗi ca lưu cả ảnh PNG và dữ liệu JSON của đường/chủ đề ROS; ảnh không phải minh họa dựng lại.", 1.75, 6.90, 9.85, 0.28, size=11.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # 32 — Geometry results
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Kết quả hình học trên 34 nhóm ghép cặp đầy đủ", "Kết quả", 32)
    add_picture_contain(slide, ASSETS / "figures" / "figure_09_aggregate_metrics.png", 0.66, 1.44, 7.74, 5.18)
    pstmo = geo["methods"]["pstmo"]
    add_metric(slide, 8.68, 1.55, 1.68, 1.52, value=f"{pstmo['paired_mean_max_abs_curvature_1pm']:.2f}", label="κ_max (m⁻¹)", accent=GREEN, note="giảm 90,3% so với Raw")
    add_metric(slide, 10.58, 1.55, 1.68, 1.52, value=f"{pstmo['paired_mean_curvature_energy_1pm']:.2f}", label="Eκ (m⁻¹)", accent=GREEN, note="giảm 98,5% so với Raw")
    add_metric(slide, 8.68, 3.38, 1.68, 1.52, value=f"{pstmo['paired_mean_path_length_m']:.2f}", label="chiều dài (m)", accent=TEAL, note="ngắn hơn Raw 1,6%")
    add_metric(slide, 10.58, 3.38, 1.68, 1.52, value=f"{pstmo['paired_mean_footprint_clearance_min_m']:.3f}", label="khoảng hở nhỏ nhất (m)", accent=ORANGE, note="không có mẫu va chạm")
    add_card(slide, 8.68, 5.25, 3.58, 1.22, title="Cách đọc đúng", body="PSTMO mượt hơn rõ rệt; khoảng hở thấp hơn một số đối chứng nhưng vẫn đạt kiểm tra hình bao.", accent=ORANGE, title_size=12.8, body_size=10.0)
    add_text(slide, "Số liệu là trung bình ghép cặp; không trộn các bộ lập kế hoạch hoặc đường gốc khác nhau.", 1.35, 6.77, 10.65, 0.28, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)

    # 33 — Success/runtime/diagnostics
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Tỷ lệ thành công, thời gian xử lý và chẩn đoán", "Kết quả", 33)
    add_picture_contain(slide, ASSETS / "figures" / "figure_10_success_points_runtime.png", 0.70, 1.46, 5.88, 3.08)
    add_picture_contain(slide, ASSETS / "figures" / "figure_11_dq_live_diagnostics.png", 6.83, 1.46, 5.80, 3.08)
    add_metric(slide, 0.84, 4.88, 2.05, 1.34, value="35/35", label="PSTMO hình học", accent=GREEN)
    add_metric(slide, 3.10, 4.88, 2.05, 1.34, value="34/35", label="Simple hình học", accent=ORANGE, note="1 ca thất bại")
    add_metric(slide, 5.36, 4.88, 2.05, 1.34, value="95,2 ms", label="thời gian xử lý PSTMO", accent=PURPLE)
    add_metric(slide, 7.62, 4.88, 2.05, 1.34, value="221", label="đoạn chuyển tiếp G²", accent=TEAL, note="tổng trên 35 đường")
    add_metric(slide, 9.88, 4.88, 2.05, 1.34, value="6", label="lần quay tại chỗ", accent=ORANGE, note="trên tổng 229 góc")
    add_text(slide, "PSTMO cần thời gian xử lý cao hơn các phương án đối chứng, đổi lại cung cấp kiểm tra và dữ liệu chẩn đoán chi tiết hơn.", 1.17, 6.58, 11.00, 0.40, size=12.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 34 — Execution results
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Kết quả thực thi vòng kín trong Gazebo", "Kết quả", 34)
    add_picture_contain(slide, ASSETS / "figures" / "figure_13_execution_overall.png", 0.68, 1.46, 6.08, 3.44)
    add_picture_contain(slide, ASSETS / "figures" / "figure_14_execution_by_planner.png", 6.98, 1.46, 5.65, 3.44)
    pair_raw = execution["pstmo_paired_comparison"]["raw"]
    pair_simple = execution["pstmo_paired_comparison"]["simple"]
    add_metric(slide, 0.86, 5.19, 2.12, 1.32, value="170/175", label="lượt tới đích", accent=GREEN, note="5 lỗi cùng nhóm C30")
    add_metric(slide, 3.22, 5.19, 2.12, 1.32, value="53,15 s", label="PSTMO trung bình", accent=GREEN)
    add_metric(slide, 5.58, 5.19, 2.12, 1.32, value=f"{abs(pair_raw['paired_difference_s_mean']):.2f} s", label="nhanh hơn Raw", accent=TEAL, note=f"{pair_raw['pstmo_faster_pair_count']}/34 cặp")
    add_metric(slide, 7.94, 5.19, 2.12, 1.32, value=f"{abs(pair_simple['paired_difference_s_mean']):.2f} s", label="nhanh hơn Simple", accent=BLUE, note=f"{pair_simple['pstmo_faster_pair_count']}/34 cặp")
    add_metric(slide, 10.30, 5.19, 2.12, 1.32, value="0", label="lần can thiệp", accent=GREEN, note="Collision Monitor")
    add_text(slide, "Thời gian được đo từ lúc FollowPath nhận lệnh tới khi robot dừng vật lý theo trạng thái thực mô phỏng.", 1.17, 6.76, 11.00, 0.30, size=11.7, color=MUTED, align=PP_ALIGN.CENTER)

    # 35 — Detailed case
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Một ca thực thi điển hình: kho điều phối · Smac2D", "Kết quả", 35)
    add_picture_contain(slide, ASSETS / "execution_cases" / "warehouse_dispatch_smac2d.png", 0.68, 1.44, 8.90, 5.55)
    add_card(slide, 9.86, 1.52, 2.82, 1.34, title="Cùng đường gốc đầu vào", body="Năm phương án dùng cùng đường do bộ lập kế hoạch tạo, xác nhận bằng hàm băm.", accent=BLUE, title_size=12.5, body_size=10.2)
    add_card(slide, 9.86, 3.08, 2.82, 1.34, title="Hai lớp bằng chứng", body="Đường kế hoạch và quỹ đạo Gazebo được đặt cạnh nhau.", accent=TEAL, title_size=13, body_size=10.8)
    add_card(slide, 9.86, 4.64, 2.82, 1.34, title="Không chỉ nhìn hình", body="Kết quả còn có thời gian, khoảng hở, va chạm và trạng thái thành công.", accent=GREEN, title_size=13, body_size=10.4)
    add_text(slide, "Mục tiêu của ảnh: kiểm tra đường được công bố khớp với dữ liệu RViz2/Gazebo, không thay thế số liệu định lượng.", 9.87, 6.32, 2.82, 0.64, size=10.8, color=MUTED, align=PP_ALIGN.CENTER)

    # 36 — C30 failure analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Phân tích ca thất bại C30", "Kết quả", 36, subtitle="Kho điều phối · SmacHybrid · không được coi thời điểm dừng sớm là thời gian hoàn thành")
    add_picture_contain(slide, ASSETS / "figures" / "case_C30_warehouse_dispatch_SmacHybrid.png", 0.62, 1.42, 7.05, 5.25)
    add_card(slide, 7.94, 1.50, 4.72, 1.35, title="Hình học của đường PSTMO", body="10 góc được xử lý: −45,77° đến +56,30°; Eκ giảm 52,84% và κmax giảm 14,77% so với đường gốc.", accent=TEAL, title_size=13, body_size=10.7)
    add_card(slide, 7.94, 3.05, 4.72, 1.35, title="Lỗi riêng của Simple", body="Máy chủ làm mượt từ chối đường tại (−4,741290; 3,482165; 0,352672 rad) do hình bao va chạm.", accent=ORANGE, title_size=13, body_size=10.7)
    add_card(slide, 7.94, 4.60, 4.72, 1.55, title="Lỗi thực thi của bốn đường còn lại", body="RPP dự báo va chạm phía trước; FollowPath kết thúc với PATIENCE_EXCEEDED, mã 104. Bộ giám sát va chạm không can thiệp: 0 lần.", accent=RED, title_size=13, body_size=10.5)
    add_text(slide, "Kết luận đúng: đây là giới hạn phối hợp lập kế hoạch–làm mượt–điều khiển, không phải bằng chứng rằng mọi đường đều va chạm hình học.", 1.03, 6.78, 11.25, 0.34, size=11.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 37 — Conclusions and limits
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Kết luận và giới hạn", "Kết luận", 37)
    add_rect(slide, 0.72, 1.48, 5.78, 4.98, fill=PALE_GREEN, line=LINE)
    add_text(slide, "KẾT LUẬN TỪ DỮ LIỆU HIỆN CÓ", 1.04, 1.82, 5.12, 0.40, size=16.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "PSTMO tạo đường G², có mô hình bánh xe, thời gian và vùng quét hình bao robot.",
        "Trên 34 nhóm ghép cặp, κ_max giảm 90,3% và Eκ giảm 98,5% so với Raw.",
        "Trong mô phỏng vòng kín, PSTMO nhanh hơn Raw trung bình 3,97 s trên các cặp thành công.",
        "Không ghi nhận mẫu va chạm trên đường kế hoạch hoặc can thiệp của bộ giám sát va chạm.",
    ], 1.06, 2.50, 5.08, 3.25, size=13.2, spacing=9)
    add_rect(slide, 6.82, 1.48, 5.78, 4.98, fill=PALE_ORANGE, line=LINE)
    add_text(slide, "GIỚI HẠN CẦN NÓI RÕ", 7.14, 1.82, 5.12, 0.40, size=16.5, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "Mỗi tổ hợp mới chạy một lần; chưa đủ để suy luận thống kê mạnh.",
        "Bản đồ tĩnh không mô hình hóa vật cản động, trượt bánh, tải và độ trễ cảm biến.",
        "Eκ là chỉ số hình học, không phải điện năng thực.",
        "Cần thử nghiệm lặp và thử robot thật với sai số bám đường, dòng điện và sự kiện an toàn.",
    ], 7.16, 2.50, 5.08, 3.25, size=13.2, spacing=9)
    add_equation(slide, "Thông điệp: PSTMO cải thiện rõ rệt chuyển hướng trong mô phỏng; robot thật là bước kiểm chứng tiếp theo.", 1.18, 6.62, 10.98, 0.58, color=NAVY, fill=LIGHT, size=14)

    # 38 — Glossary
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_title(slide, "Thuật ngữ nên dùng nhất quán", "Phụ lục", 38, subtitle="Thuật ngữ tiếng Anh · thuật ngữ tiếng Việt · cách hiểu trong bài")
    data = [
        ["English term", "Tiếng Việt nên dùng", "Ý nghĩa trong báo cáo"],
        ["path", "đường đi / đường hình học", "Chuỗi tư thế chưa gắn thời gian"],
        ["trajectory", "quỹ đạo theo thời gian", "Có t, v, ω, gia tốc"],
        ["path smoother", "bộ làm mượt đường đi", "Chỉnh hình học đường đi"],
        ["corner-handling option", "phương án xử lý tại góc", "Một cách xử lý cụ thể cho một góc"],
        ["trim distance d", "khoảng cắt d", "Chiều dài cắt trên mỗi cạnh"],
        ["shape ratio α=q/d", "tỷ lệ hình dạng α", "Điều chỉnh phân bố độ cong"],
        ["curvature κ", "độ cong κ", "Mức đổi hướng trên một đơn vị chiều dài"],
        ["curvature energy Eκ", "tích phân bình phương độ cong", "Chỉ số uốn hình học; không phải điện năng"],
        ["speed limit / speed profile", "giới hạn vận tốc / biểu đồ vận tốc", "Ngưỡng cục bộ / dãy v(s) sau quét"],
        ["swept-footprint", "vùng quét hình bao robot", "Hợp của footprint dọc chuyển động"],
        ["time gate", "điều kiện ưu thế thời gian", "Giữ đoạn chuyển tiếp khi có lợi hơn quay tại chỗ"],
        ["pivot", "quay tại chỗ", "Hai bánh quay ngược chiều để đổi góc hướng"],
    ]
    add_table(slide, data, 0.64, 1.44, 12.05, 5.78, col_widths=[3.10, 3.55, 5.40], font_size=9.6)

    # 39 — Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_text(slide, "CẢM ƠN THẦY CÔ VÀ CÁC BẠN", 1.15, 1.22, 11.05, 0.46, size=17, color="B8C8D8", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Q & A", 1.15, 2.20, 11.05, 1.00, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "PSTMO · ROS 2 Navigation2 · Bézier bậc năm G²", 1.15, 3.46, 11.05, 0.42, size=17, color="A9D7CC", align=PP_ALIGN.CENTER)
    add_equation(slide, "Đường mượt hơn chỉ có ý nghĩa khi vẫn khả thi, an toàn và thực thi được.", 2.05, 4.50, 9.25, 0.74, color=NAVY, fill=WHITE, size=18)
    add_text(slide, "NGUYỄN TIẾN CƯƠNG · 2026", 1.15, 6.40, 11.05, 0.34, size=12, color="91A4B8", align=PP_ALIGN.CENTER)

    SLIDE_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
